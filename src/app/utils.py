import os
import sys
import uuid
import tempfile 
import io
import logging
from datetime import datetime
from typing import List, Optional, Dict, Any, Union
from pathlib import Path
import pypdf
import docx2python
from pptx import Presentation
from openpyxl import load_workbook

# Handle magic import for Windows vs other OS
if sys.platform == 'win32':
    import magic  # This will use python-magic-bin on Windows
else:
    import magic  # This will use python-magic on other platforms

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# Configure logging
logger = logging.getLogger(__name__)

# Supported file types and their MIME types
SUPPORTED_FILE_TYPES = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/msword': '.doc',
    'application/vnd.ms-powerpoint': '.ppt',
    'application/vnd.ms-excel': '.xls',
    'text/plain': '.txt'
}

def get_file_extension(file_path: str) -> Optional[str]:
    """Get file extension based on MIME type."""
    mime = magic.Magic(mime=True)
    file_mime = mime.from_file(file_path)
    return SUPPORTED_FILE_TYPES.get(file_mime)

def save_upload_file(upload_file, suffix: str = "") -> str:
    """Save uploaded file to a temporary file and return its path."""
    try:
        suffix = "." + suffix.lstrip(".") if suffix else ""
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(upload_file.file.read())
            return tmp.name
    except Exception as e:   
        logger.error(f"Error saving uploaded file: {str(e)}")
        raise 

def load_document(file_path: str) -> List[Document]:
    """Load document using appropriate loader based on file type."""
    try:
        file_ext = get_file_extension(file_path)
        
        if not file_ext: 
            # Try to get the extension from the filename as a fallback
            file_ext = os.path.splitext(file_path)[1].lower()
            if not file_ext or file_ext not in ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.txt']:
                raise ValueError(f"Unsupported file type: {file_path}. Supported types: {', '.join(SUPPORTED_FILE_TYPES.values())}")
        
        logger.info(f"Loading document with extension: {file_ext}")
        
        if file_ext == '.pdf':
            # Load PDF using pypdf
            reader = pypdf.PdfReader(file_path)
            text = "\n\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif file_ext in ['.docx', '.doc']:
            # Load DOCX using docx2python
            doc = docx2python.docx2python(file_path)
            text = doc.text
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif file_ext in ['.pptx', '.ppt']:
            # Load PPTX using python-pptx
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif file_ext in ['.xlsx', '.xls']:
            # Load XLSX using openpyxl
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join(str(cell) for cell in row if cell is not None) + "\n"
            return [Document(page_content=text, metadata={"source": file_path})]
            
        else:
            # For .txt files or any other text files
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return [Document(page_content=text, metadata={"source": file_path})]
            
    except Exception as e:
        logger.error(f"Error loading document {os.path.basename(file_path)}: {str(e)}", exc_info=True)
        raise

def chunk_documents(documents: List[Document], 
                   chunk_size: int = 1000, 
                   chunk_overlap: int = 200) -> List[Document]:
    """Split documents into chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    return text_splitter.split_documents(documents)

def process_uploaded_file(upload_file, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Document]:
    """Process an uploaded file: save, load, and chunk it."""
    temp_file_path = None
    try:
        if not upload_file or upload_file.filename == '':
            raise ValueError("No file was uploaded")
            
        logger.info(f"Processing uploaded file: {upload_file.filename} (size: {upload_file.size} bytes)")
        
        # Save the uploaded file temporarily with its original extension
        file_ext = os.path.splitext(upload_file.filename)[1]
        temp_file_path = save_upload_file(upload_file, suffix=file_ext)
        
        # Verify the file has content
        file_size = os.path.getsize(temp_file_path)
        if file_size == 0:
            raise ValueError("Uploaded file is empty")
            
        logger.info(f"Temporary file saved: {temp_file_path} (size: {file_size} bytes)")
        
        # Load the document
        documents = load_document(temp_file_path)
        if not documents:
            raise ValueError("No content could be extracted from the file")
        
        logger.info(f"Successfully loaded {len(documents)} document(s)")
        
        # Chunk the document
        chunks = chunk_documents(documents, chunk_size, chunk_overlap)
        if not chunks:
            raise ValueError("No chunks could be created from the document")
            
        logger.info(f"Created {len(chunks)} chunks from the document")
        return chunks
        
    except Exception as e:
        error_msg = f"Error processing file {getattr(upload_file, 'filename', 'unknown')}: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise ValueError(error_msg) from e
        
    finally:
        # Clean up the temporary file
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
                logger.debug(f"Temporary file deleted: {temp_file_path}")
            except Exception as e:
                logger.error(f"Error deleting temporary file {temp_file_path}: {str(e)}")

def generate_slide_metadata(slide: Optional[dict] = None) -> dict:
    """Generate metadata dictionary for a slide."""
    now = datetime.utcnow().isoformat()
    if slide is None:
        return {
            "slide_name": "",
            "course_name": "",
            "subject_name": "",
            "description": "",
            "created_at": now,
            "updated_at": now
        }
    return {
        "slide_name": getattr(slide, 'slide_name', ''),
        "course_name": getattr(slide, 'course_name', ''),
        "subject_name": getattr(slide, 'subject_name', ''),
        "description": getattr(slide, 'description', '') or "",
        "created_at": now,
        "updated_at": now
    }

def generate_slide_id() -> str:
    """Generate a unique ID for a slide."""
    return str(uuid.uuid4())

def extract_text_from_pdf(file_content: bytes) -> str:
    """Legacy function for backward compatibility."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        
        pdf_reader = PyPDF2.PdfReader(tmp_path)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)